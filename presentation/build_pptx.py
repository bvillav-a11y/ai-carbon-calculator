#!/usr/bin/env python3
"""Generate an editable, Google-Slides-importable .pptx of the AI Carbon Calculator talk.

Mirrors presentation/index.html (the reveal.js deck): same dark theme + app design
tokens, the three pillars, and per-slide speaker notes (kept in each slide's notes pane,
so they survive the import into Google Slides).

Everything is native PowerPoint text boxes / shapes — fully editable after import.
Fonts are named "Space Grotesk" / "Space Mono"; enable them in Google Slides via
"More fonts" for an exact match (otherwise Slides substitutes a default).

Regenerate:  python3 presentation/build_pptx.py
Output:      presentation/ai-carbon-calculator.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── App design tokens (from index.html :root, dark theme) ──
BG       = RGBColor(0x0d, 0x11, 0x17)
SURFACE  = RGBColor(0x16, 0x1b, 0x22)
SURFACE2 = RGBColor(0x1c, 0x23, 0x30)
BORDER   = RGBColor(0x44, 0x4c, 0x56)
GREEN    = RGBColor(0x3f, 0xb9, 0x50)
GREENDIM = RGBColor(0x12, 0x2d, 0x1a)
TEXT     = RGBColor(0xe6, 0xed, 0xf3)
MUTED    = RGBColor(0xb1, 0xba, 0xc4)
DIM      = RGBColor(0x8b, 0x94, 0x9e)
BLUE     = RGBColor(0x79, 0xb8, 0xff)
PURPLE   = RGBColor(0xbc, 0x8c, 0xff)
TEAL     = RGBColor(0x2f, 0xd4, 0xc6)
ORANGE   = RGBColor(0xff, 0x9d, 0x3c)
PINK     = RGBColor(0xff, 0x8a, 0xd4)

SANS = "Space Grotesk"
MONO = "Space Mono"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
MARGIN = Inches(0.9)
CONTENT_W = prs.slide_width - 2 * MARGIN

# ── run-spec helpers: (text, style-dict) ──
def S(t, sz=18, c=TEXT, b=False, i=False, f=SANS):
    return (t, dict(font=f, size=sz, color=c, bold=b, italic=i))

def M(t, sz=15, c=GREEN, b=False):
    return (t, dict(font=MONO, size=sz, color=c, bold=b))

def _style_run(r, text, font=SANS, size=18, color=TEXT, bold=False, italic=False):
    r.text = text
    f = r.font
    f.name = font
    f.size = Pt(size)
    f.color.rgb = color
    f.bold = bold
    f.italic = italic

def _fill_para(p, runs):
    for (t, kw) in runs:
        _style_run(p.add_run(), t, **kw)

# ── primitives ──
def new_slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def text(slide, left, top, w, h, paragraphs, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP, line_spacing=None, space_after=None):
    """paragraphs: list of run-lists; each run-list is [(text, style-dict), ...]."""
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after is not None:
            p.space_after = Pt(space_after)
        _fill_para(p, runs)
    return tb

def no_shadow(shp):
    shp.shadow.inherit = False

def rrect(slide, left, top, w, h, fill=SURFACE, line=BORDER, line_w=1.0, radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, int(w), int(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    no_shadow(shp)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    return shp

def label(slide, t, top=Inches(0.7)):
    text(slide, MARGIN, top, CONTENT_W, Inches(0.4),
         [[M(t.upper(), sz=14, c=GREEN, b=True)]])

def title(slide, t, top=Inches(1.18), size=40):
    text(slide, MARGIN, top, CONTENT_W, Inches(1.25),
         [[S(t, sz=size, c=TEXT, b=True)]], line_spacing=1.0)

def footnote(slide, runs, top=Inches(6.45)):
    text(slide, MARGIN, top, CONTENT_W, Inches(0.7), [runs], line_spacing=1.1)

def notes(slide, t):
    slide.notes_slide.notes_text_frame.text = t

def pill(slide, t, top, left=MARGIN):
    w = Inches(0.155 * len(t) + 0.9)
    shp = rrect(slide, left, top, w, Inches(0.5), fill=GREENDIM, line=None, radius=0.5)
    tf = shp.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _fill_para(p, [M(t.upper(), sz=13, c=GREEN, b=True)])
    return shp

def card(slide, left, top, w, h, key, body_paras):
    shp = rrect(slide, left, top, w, h, fill=SURFACE)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.15); tf.margin_bottom = Inches(0.15)
    p0 = tf.paragraphs[0]
    _fill_para(p0, [M(key.upper(), sz=12, c=GREEN, b=True)])
    p0.space_after = Pt(7)
    for runs in body_paras:
        p = tf.add_paragraph()
        p.line_spacing = 1.12
        _fill_para(p, runs)
    return shp

def cards_row(slide, items, top, height, gap=Inches(0.4)):
    n = len(items)
    col_w = int((CONTENT_W - gap * (n - 1)) / n)
    x = MARGIN
    for key, body in items:
        card(slide, x, top, col_w, height, key, body)
        x += col_w + gap

def cards_stack(slide, items, top, item_h, gap=Inches(0.25)):
    y = top
    for key, body in items:
        card(slide, MARGIN, int(y), CONTENT_W, item_h, key, body)
        y += item_h + gap

def bullets(slide, items, top=Inches(2.2), size=20, gap=11):
    tb = slide.shapes.add_textbox(MARGIN, top, CONTENT_W, Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, runs in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.12
        _style_run(p.add_run(), "•  ", font=SANS, size=size, color=GREEN, bold=True)
        _fill_para(p, runs)
    return tb

def ledger(slide, rows, top, row_h=Inches(0.52), pad=Inches(0.2), lw=0.62):
    n = len(rows)
    H = int(row_h * n + pad * 2)
    rrect(slide, MARGIN, top, CONTENT_W, H, fill=SURFACE2)
    y = top + pad
    for left_runs, val_runs in rows:
        lt = slide.shapes.add_textbox(MARGIN + pad, int(y), int(CONTENT_W * lw), row_h)
        lt.text_frame.word_wrap = True
        lt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _fill_para(lt.text_frame.paragraphs[0], left_runs)
        vt = slide.shapes.add_textbox(MARGIN + int(CONTENT_W * lw), int(y),
                                      int(CONTENT_W * (1 - lw)) - pad, row_h)
        vt.text_frame.word_wrap = True
        vt.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        pv = vt.text_frame.paragraphs[0]
        pv.alignment = PP_ALIGN.RIGHT
        _fill_para(pv, val_runs)
        y += row_h

def codeblock(slide, lines, top=Inches(2.4), height=Inches(1.9)):
    shp = rrect(slide, MARGIN, top, CONTENT_W, height, fill=SURFACE2)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.22); tf.margin_top = Inches(0.18)
    for i, runs in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        _fill_para(p, runs)

def flow(slide, nodes, top=Inches(2.7), gap=Inches(0.5), h=Inches(1.0)):
    n = len(nodes)
    node_w = int((CONTENT_W - gap * (n - 1)) / n)
    x = MARGIN
    for i, t in enumerate(nodes):
        shp = rrect(slide, x, top, node_w, h, fill=SURFACE)
        shp.text_frame.word_wrap = True
        shp.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = shp.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _fill_para(p, [M(t, sz=14, c=TEXT)])
        if i < n - 1:
            ar = slide.shapes.add_textbox(x + node_w, top, gap, h)
            ar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            pa = ar.text_frame.paragraphs[0]
            pa.alignment = PP_ALIGN.CENTER
            _fill_para(pa, [S("→", sz=24, c=GREEN, b=True)])
        x += node_w + gap

def eqgrid(slide, items, top=Inches(2.8), gap=Inches(0.5)):
    n = len(items)
    col_w = int((CONTENT_W - gap * (n - 1)) / n)
    x = MARGIN
    for num, numcolor, cap in items:
        text(slide, x, top, col_w, Inches(1.2), [[M(num, sz=46, c=numcolor, b=True)]],
             align=PP_ALIGN.CENTER)
        text(slide, x, top + Inches(1.35), col_w, Inches(0.8), [[M(cap, sz=14, c=MUTED)]],
             align=PP_ALIGN.CENTER)
        x += col_w + gap

def confbar(slide, top, lab, level, pct):
    text(slide, MARGIN, top, CONTENT_W, Inches(0.34), [[M(lab, sz=15, c=MUTED)]])
    text(slide, MARGIN, top, CONTENT_W, Inches(0.34), [[M(level, sz=15, c=MUTED)]],
         align=PP_ALIGN.RIGHT)
    rrect(slide, MARGIN, top + Inches(0.4), CONTENT_W, Inches(0.22), fill=SURFACE2,
          line=BORDER, line_w=0.75, radius=0.5)
    rrect(slide, MARGIN, top + Inches(0.4), int(CONTENT_W * pct), Inches(0.22),
          fill=GREEN, line=None, radius=0.5)

def divider(pilltext, titletext):
    s = new_slide()
    pill(s, pilltext, top=Inches(2.7))
    text(s, MARGIN, Inches(3.45), CONTENT_W, Inches(1.6),
         [[S(titletext, sz=48, c=TEXT, b=True)]], line_spacing=1.0)
    return s


# ════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════

# 1 · TITLE
s = new_slide()
pill(s, "AI Carbon Calculator", top=Inches(2.0))
text(s, MARGIN, Inches(2.7), CONTENT_W, Inches(1.8),
     [[S("What does your AI use", sz=50, c=TEXT, b=True)],
      [S("cost the planet?", sz=50, c=TEXT, b=True)]], line_spacing=1.0)
text(s, MARGIN, Inches(4.7), CONTENT_W, Inches(0.6),
     [[S("The logic, the bases, and the deployment of a footprint estimator.", sz=20, c=MUTED)]])
text(s, MARGIN, Inches(5.4), CONTENT_W, Inches(0.5),
     [[M("Bianca Villavicencio · TELUS Digital · 2026-06-16", sz=14, c=MUTED)]])
notes(s, "Title. One line: \"This is a tour of a small tool that estimates the monthly "
         "carbon footprint of personal AI use — and more interestingly, the reasoning and "
         "the deployment story behind it.\" Set the frame: three parts — the choices, the "
         "bases (the math + data), and getting it live on gizmos. ~30s.")

# 2 · THE QUESTION
s = new_slide()
label(s, "the problem")
title(s, "Nobody actually knows their number")
bullets(s, [
    [S("Providers don't disclose per-query energy, water, or training amortisation.", 20)],
    [S("\"AI is bad for the environment\" is a vibe, not a number you can act on.", 20)],
    [S("Goal: a ", 20), S("defensible estimate", 20, b=True),
     S(" from inputs a normal person can answer — and show the working.", 20)],
], top=Inches(2.3))
footnote(s, [M("Live: ai-carbon-footprint.telus.gizmos.run  ·  public mirror: "
               "bvillav-a11y.github.io/ai-carbon-calculator", sz=13, c=MUTED)])
notes(s, "The motivation. Emphasise the gap: people feel guilty or dismissive, neither is "
         "grounded. The bar we set ourselves: defensible, transparent, and driven by answers "
         "a non-expert can give. Mention it's live on gizmos internally and mirrored on "
         "GitHub Pages. ~45s.")

# DIVIDER · PILLAR I
s = divider("Pillar I", "The logic of the choices")
notes(s, "Transition. \"First, the product reasoning — why the tool is shaped the way it is.\" ~10s.")

# 3 · MENTAL MODEL
s = new_slide()
label(s, "the mental model")
title(s, "One pipeline, four stages")
flow(s, ["survey\nanswers", "inferred\nparameters", "carbon\nmath", "tangible\nequivalents"])
footnote(s, [S("Everything downstream is a pure function of the inferred parameters, so any "
               "value can be overridden and the whole chain re-renders live.", 15, c=MUTED)],
         top=Inches(4.3))
notes(s, "This is the spine of the whole app — keep coming back to it. Plain-language answers "
         "become physical parameters; parameters feed one math function; the kg result is "
         "translated into things people feel. The key engineering property: it's a one-way "
         "pure pipeline, so \"what-if\" is free — change an input, everything recomputes. ~50s.")

# 4 · WHY A SURVEY
s = new_slide()
label(s, "choice · inputs")
title(s, "A survey, not a parameter form")
cards_row(s, [
    ("you ask", [[S("\"How complex are your tasks?\"", 16)],
                 [S("\"Cloud or on-prem?\"", 16)],
                 [S("\"Interactive or batch?\"", 16)]]),
    ("you infer", [[S("task complexity → ", 16), S("model tier", 16, b=True)],
                   [S("hardware type → ", 16), S("PUE", 16, b=True)],
                   [S("interactivity + frequency → ", 16), S("utilisation", 16, b=True)]]),
], top=Inches(2.3), height=Inches(2.2))
footnote(s, [S("Nobody knows their PUE. Everybody knows whether they run batch jobs or chat.",
               16, c=MUTED)], top=Inches(4.9))
notes(s, "The core product bet: meet people where they are. Asking for PUE or utilisation "
         "directly would lose everyone. So we map answerable questions onto the physical "
         "parameters. Note the deliberate decoupling: task complexity drives the model tier; "
         "interactivity feeds utilisation, NOT the tier. Each inference also stores a reason "
         "string, surfaced as a tooltip on the matching control. ~55s.")

# 5 · OUTCOME-FIRST UX
s = new_slide()
label(s, "choice · layout")
title(s, "Outcome first, knobs last")
bullets(s, [
    [S("Results", 20, b=True), S(" — the number + tangible equivalents, immediately.", 20)],
    [S("Insights & actions", 20, b=True),
     S(" — right-size your model, cleaner grid, concrete cuts.", 20)],
    [S("More information", 20, b=True),
     S(" — parameters, calculation trace, confidence, sources — all collapsed.", 20)],
], top=Inches(2.3))
footnote(s, [S("Most people want the answer and what to do about it. Experts can open every knob.",
               16, c=MUTED)])
notes(s, "Progressive disclosure. The old design led with a parameter form; we flipped it. "
         "Land on the outcome, then \"here's what you can do,\" then — for the curious — the "
         "full machinery behind four collapsibles under a \"More information\" header. This is "
         "also where engagement telemetry lives: we record which of those menus people open. ~50s.")

# 6 · CREDIBILITY BY DESIGN
s = new_slide()
label(s, "choice · trust")
title(s, "Credibility is a feature")
cards_row(s, [
    ("every assumption", [[S("editable ", 16), S("and", 16, i=True),
                           S(" sourced — with a confidence bar per factor.", 16)]]),
    ("honest by default", [[S("\"This is probably an underestimate\" — we say what's excluded.",
                              16)]]),
], top=Inches(2.3), height=Inches(1.7))
footnote(s, [S("Reserved impact-colour (neutral scenario tiers)  ·  WCAG 2.1 AA  ·  "
               "ARIA radio-groups + live regions.", 16, c=MUTED)], top=Inches(4.5))
notes(s, "For a sceptical technical audience, trust is the product. Three moves: (1) every "
         "number is both editable and cited, with an explicit confidence level; (2) we "
         "proactively state the estimate is conservative and list what's left out; "
         "(3) accessibility is maintained, not bolted on — ARIA patterns, keyboard nav, AA "
         "contrast (several palette colours were lifted specifically to pass AA). ~55s.")

# DIVIDER · PILLAR II
s = divider("Pillar II", "The bases: the math & the data")
notes(s, "Transition into the rigour. \"Now the part a technical audience actually wants — "
         "what's under the number.\" ~10s.")

# 7 · THE CALCULATION CHAIN
s = new_slide()
label(s, "computeCarbon()")
title(s, "The chain, as a ledger")
ledger(s, [
    ([M("raw energy = inTok·in_kWh/kTok + outTok·out_kWh/kTok", sz=13, c=TEXT)], [M("kWh", sz=14)]),
    ([M("÷ utilisation", sz=14, c=TEXT)], [M("kWh", sz=14)]),
    ([M("× PUE", sz=14, c=TEXT)], [M("kWh", sz=14)]),
    ([M("× grid g/kWh", sz=14, c=TEXT)], [M("operational kg", sz=14)]),
    ([M("+ embodied  (GPU manufacture, amortised)", sz=14, c=TEXT)], [M("kg", sz=14)]),
    ([M("= total kg CO₂", sz=14, c=TEXT, b=True)], [M("total", sz=14)]),
], top=Inches(2.1))
footnote(s, [M("One pure function. The on-screen \"calculation trace\" renders exactly these rows.",
               sz=13, c=MUTED)])
notes(s, "Walk each line slowly — this is the meatiest slide. Energy starts per-token, split "
         "input vs output. Divide by utilisation (idle GPUs still draw power). Multiply by PUE "
         "(data-centre overhead). Multiply by grid intensity to turn kWh into kg. Add amortised "
         "embodied carbon. The on-screen ledger in the app mirrors this row for row — the math "
         "is never hidden. ~90s; this slide can run long, that's fine.")

# 8 · ENERGY PER TOKEN
s = new_slide()
label(s, "MODEL_SPECS · kWh per 1k tokens")
title(s, "Output costs ~5× input")
ledger(s, [
    ([M("Claude Haiku", sz=15, c=TEXT)],
     [M("in 0.000090", sz=15, c=BLUE), M("  ·  ", sz=15, c=MUTED), M("out 0.000450", sz=15, c=PURPLE)]),
    ([M("Claude Sonnet", sz=15, c=TEXT)],
     [M("in 0.000390", sz=15, c=BLUE), M("  ·  ", sz=15, c=MUTED), M("out 0.001950", sz=15, c=PURPLE)]),
    ([M("Claude Opus", sz=15, c=TEXT)],
     [M("in 0.000780", sz=15, c=BLUE), M("  ·  ", sz=15, c=MUTED), M("out 0.003900", sz=15, c=PURPLE)]),
], top=Inches(2.1), lw=0.5)
text(s, MARGIN, Inches(4.35), CONTENT_W, Inches(0.5),
     [[S("■ input", 16, c=BLUE), S(" — cheap, parallel    ", 16, c=MUTED),
       S("■ output", 16, c=PURPLE), S(" — generated one token at a time, ~5× costlier", 16, c=MUTED)]])
footnote(s, [M("Source: Epoch AI — \"How much energy does ChatGPT use?\" (Sonnet ≈ GPT-4o class).",
               sz=13, c=MUTED)])
notes(s, "Two points. First, energy is per-token and asymmetric — output is autoregressive (one "
         "token at a time) so it's roughly 5× input; that's why \"ask for concise answers\" is a "
         "real lever. Second, tiers scale ~linearly Haiku→Sonnet→Opus, which is what makes "
         "\"right-size your model\" a genuine carbon decision, not a rounding error. Numbers are "
         "grounded in Epoch AI's public estimates. ~60s.")

# 9 · HARDWARE & OVERHEAD
s = new_slide()
label(s, "GPU_SPECS · PUE · utilisation")
title(s, "Where the watts leak")
cards_stack(s, [
    ("gpu", [[S("H100 700W  ·  A100 400W (SXM4) / 300W (PCIe)  ·  embodied 130–225 kg CO₂ per card",
               16)]]),
    ("pue", [[S("1.2 hyperscaler  ·  1.4 typical cloud  ·  2.0+ on-prem — multiplies IT energy "
               "by overhead", 16)]]),
    ("utilisation", [[S("energy per ", 16), S("useful", 16, i=True),
                      S(" token scales as ", 16), S("1 ÷ utilisation", 16, b=True),
                      S(" — idle silicon still burns", 16)]]),
], top=Inches(2.2), item_h=Inches(1.15))
notes(s, "Three multipliers people never see. GPU choice sets the power draw and carries "
         "embodied carbon. PUE is the tax for cooling and networking — 1.4 means 40% on top. "
         "Utilisation is the sneaky one: a half-idle GPU doubles the energy attributed to each "
         "real token, which is exactly why interactive use is dirtier per token than batch. ~60s.")

# 10 · THE GRID
s = new_slide()
label(s, "GRID_MAP · g CO₂ / kWh")
title(s, "The grid often decides it")
ledger(s, [
    ([M("🇨🇦 Canada BC (hydro)", sz=15, c=TEXT)], [M("15", sz=15)]),
    ([M("🇫🇷 France (nuclear)", sz=15, c=TEXT)], [M("60", sz=15)]),
    ([M("🇧🇷 Brazil", sz=15, c=TEXT)], [M("136", sz=15)]),
    ([M("🇺🇸 US average", sz=15, c=TEXT)], [M("400", sz=15)]),
    ([M("🇵🇱 Poland (coal)", sz=15, c=TEXT)], [M("700", sz=15)]),
], top=Inches(2.1))
footnote(s, [S("~47× spread, end to end. Insight fires when grid > 150: \"a cleaner grid (~50) "
               "would cut X%.\"", 15, c=MUTED)])
notes(s, "The single biggest swing in the whole model. Same tokens, same hardware: a BC-hydro "
         "grid versus Polish coal differ by ~47×. So the app's biggest lever isn't the model — "
         "it's where the data centre is. That's why one of the personalised insights is "
         "grid-based: above 150 g we compute the saving against a ~50 g hydro/nuclear reference. ~55s.")

# 11 · EMBODIED CARBON
s = new_slide()
label(s, "embodied")
title(s, "The carbon you already spent")
codeblock(s, [
    [M("embodied_kg = gpuCo2 / (gpuTok · gpuLife · 12) · tokens", sz=15, c=TEXT)],
    [M("// e.g. 150 kg / (50M · 4yr · 12mo) · your tokens", sz=14, c=DIM)],
], top=Inches(2.3), height=Inches(1.4))
footnote(s, [S("Manufacturing CO₂ amortised over the GPU's lifetime token throughput. Small next "
               "to operational — but counted, because honesty.", 16, c=MUTED)], top=Inches(4.1))
notes(s, "Most footprint tools ignore embodied carbon entirely. We amortise the chip's "
         "manufacturing emissions across all the tokens it will ever serve, then attribute your "
         "slice. It's usually small versus operational energy — but including it is part of "
         "being defensible rather than convenient, and it's adjustable in the advanced panel. ~45s.")

# 12 · KG TO MEANING
s = new_slide()
label(s, "tangible equivalents")
title(s, "From kg to something you feel")
eqgrid(s, [
    ("6 km", ORANGE, "driven per kg CO₂"),
    ("22 kg", TEAL, "absorbed / tree / year"),
    ("$50", PINK, "social cost / tonne"),
])
footnote(s, [S("Each equivalent gets its own accent colour — deliberately distinct from every "
               "semantic colour in the math.", 16, c=MUTED)])
notes(s, "A kilogram of CO₂ is meaningless to most people. So we translate: kilometres driven, "
         "trees-years to offset, dollars of social cost. Note the colour discipline — these "
         "three accents (orange/teal/pink) are reserved exclusively for equivalents, never "
         "reused for a parameter, so the eye learns \"this colour = a real-world analogy.\" ~45s.")

# 13 · HONESTY ABOUT UNCERTAINTY
s = new_slide()
label(s, "confidence & what's excluded")
title(s, "Say what you don't know")
confbar(s, Inches(2.3), "Energy per token", "med", 0.60)
confbar(s, Inches(3.2), "PUE (inferred, undisclosed)", "low–med", 0.55)
confbar(s, Inches(4.1), "Grid intensity", "high", 0.80)
footnote(s, [S("Excluded: training amortisation · networking · cooling water · non-GPU embodied. "
               "computeCarbon is one pure function — the scenario tiers just call it with a "
               "different model.", 15, c=MUTED)], top=Inches(5.2))
notes(s, "Close the rigour section on humility. Each factor carries a confidence level — grid is "
         "well-measured, PUE is inferred and uncertain. We list the big exclusions explicitly so "
         "nobody thinks the number is complete. And the engineering punchline: because "
         "computeCarbon is a single pure function, the Haiku/Sonnet/Opus scenario comparison is "
         "literally the same function called three times — no duplicated math to drift. ~55s.")

# DIVIDER · PILLAR III
s = divider("Pillar III", "The process of deployment")
notes(s, "Transition to the war-stories. \"Now how it actually got live — and everything that "
         "fought back.\" ~10s.")

# 14 · TWO HOMES
s = new_slide()
label(s, "where it lives")
title(s, "Two targets, one source file")
cards_row(s, [
    ("GitHub Pages", [[S("public demo · serves index.html directly · /collect 404s, collection "
                         "no-ops", 16)]]),
    ("TELUS gizmos", [[S("the real home · same Worker serves the app ", 16), S("and", 16, i=True),
                       S(" a same-origin POST /collect → per-app D1", 16)]]),
], top=Inches(2.3), height=Inches(1.9))
footnote(s, [M("The entire calculator is one index.html — HTML + one CSS block + one vanilla JS "
               "block. No framework.", sz=13, c=MUTED)], top=Inches(4.6))
notes(s, "Set up the deployment shape. The calculator is a single no-framework index.html, which "
         "makes the public mirror trivial. The interesting target is gizmos — TELUS's internal "
         "app platform — where the same file is served by a Worker that also collects anonymous "
         "survey data into a per-app SQLite (D1) database. The constraint that everything is one "
         "HTML file drives the whole embed story coming up. ~50s.")

# 15 · GIZMOS ARCHITECTURE
s = new_slide()
label(s, "gizmos · architecture")
title(s, "One plain Worker does everything")
codeblock(s, [
    [M("GET  /         → serve base64-embedded index.html", sz=15, c=TEXT)],
    [M("POST /collect  → upsert survey row into D1 (by sessionId)", sz=15, c=TEXT)],
    [M("GET  /export   → CSV dump of responses", sz=15, c=TEXT)],
], top=Inches(2.3), height=Inches(1.7))
footnote(s, [S("Per-app D1 auto-provisioned · table created lazily (CREATE TABLE IF NOT EXISTS) · "
               "consent-gated.", 16, c=MUTED)], top=Inches(4.4))
notes(s, "The shipped architecture, deliberately minimal. One Worker, three routes. The HTML is "
         "embedded in the Worker as a base64 string. Collection is a same-origin POST, upserting "
         "by session id so re-runs update one row. Export is a CSV endpoint we use to verify "
         "data, because the platform's DB CLI needs a browser session. Now — why does it look "
         "like this? Because three things broke first. ~50s.")

# 16 · WAR-STORIES
s = new_slide()
label(s, "what fought back")
title(s, "Three platform lessons")
cards_stack(s, [
    ("1 · module load", [[S("a bare import 'hono' → ", 15), S("500 on load", 15, b=True),
                          S(", empty logs → rewrite as a plain Workers module, no imports", 15)]]),
    ("2 · reserved path", [[S("/api/* is platform-reserved → app POSTs ", 15),
                            S("silently blocked", 15, b=True), S(" → move to /collect & /export", 15)]]),
    ("3 · cross-origin auth", [[S("Apps-Script POSTs ", 15), S("dropped the auth cookie", 15, b=True),
                                S(" (3rd-party-cookie stripping) → pivot to ", 15),
                                S("same-origin", 15, i=True), S(" collection", 15)]]),
], top=Inches(2.2), item_h=Inches(1.15))
notes(s, "The fun slide — tell these as stories. (1) The deploy 500'd with empty logs; bisection "
         "showed the loader won't resolve a bare hono import at runtime, so any framework worker "
         "dies on load — fix was a plain Workers module with zero imports. (2) Collection POSTs "
         "vanished; turns out /api/* is reserved by the platform and app POSTs to it are blocked "
         "— renaming to /collect fixed it instantly. (3) The original plan POSTed to a Google "
         "Apps Script, but background cross-origin requests get their auth cookie stripped, so it "
         "silently 401'd — which is the whole reason collection became same-origin. Each failure "
         "shaped the architecture on the last two slides. ~90s.")

# 17 · EMBED PIPELINE
s = new_slide()
label(s, "build · embed.mjs")
title(s, "One source file, regenerated")
codeblock(s, [
    [M("index.html  +  worker.template.ts", sz=15, c=TEXT)],
    [M("      │  node scripts/embed.mjs   ", sz=15, c=TEXT), M("// base64-inline the HTML", sz=14, c=DIM)],
    [M("      ▼", sz=15, c=TEXT)],
    [M("gizmos-app/src/index.ts  →  gizmos push --org telus", sz=15, c=TEXT)],
], top=Inches(2.3), height=Inches(2.0))
footnote(s, [S("Base64, not a template literal — the inline <script> is full of backticks and "
               "${...}. Regenerate after every edit.", 16, c=MUTED)], top=Inches(4.6))
notes(s, "How a change ships. index.html stays the single source of truth; a small build step "
         "base64-encodes it into the Worker template, producing the deployable file we push. Why "
         "base64 and not a JS template string? The app's inline script is full of backticks and "
         "dollar-brace — it would shred a template literal. The one discipline: re-run embed "
         "after any HTML change, or you deploy stale. ~45s.")

# 18 · COLLECTION & TELEMETRY
s = new_slide()
label(s, "collection · telemetry")
title(s, "Anonymous, consent-gated, one row per visit")
bullets(s, [
    [S("Upsert by ", 20), M("sessionId", 18, GREEN),
     S("; baseline frozen at first results render.", 20)],
    [S("Engagement: ", 20), M("dwell_ms", 18, GREEN),
     S(", which menus opened, share, retake.", 20)],
    [S("Schema evolves via best-effort ", 20), M("ALTER TABLE … ADD COLUMN", 18, GREEN),
     S(" on an existing table.", 20)],
], top=Inches(2.3))
notes(s, "What we learn and how. Only anonymous signals — the survey answers, the computed "
         "number, and engagement: how long on results, which \"More information\" menus they "
         "opened, whether they shared or retook. All consent-gated, one upserted row per "
         "session. The schema note matters for the next slide's lesson: the prod table predates "
         "the telemetry columns, so the Worker runs idempotent ALTER TABLE migrations to add "
         "them — CREATE IF NOT EXISTS won't touch an existing table. ~55s.")

# 19 · TAKEAWAYS
s = new_slide()
label(s, "what to steal")
title(s, "Four things that generalise")
bullets(s, [
    [S("Platform constraints shape architecture", 19, b=True),
     S(" — the plain Worker + same-origin collection are scars, not preferences.", 19)],
    [S("Ship the schema and its migration together", 19, b=True),
     S(" — CREATE IF NOT EXISTS never alters a live table.", 19)],
    [S("Transparency is a feature", 19, b=True),
     S(" — editable, sourced, confidence-rated beats a confident black box.", 19)],
    [S("Right-size your model", 19, b=True),
     S(" — the app's own advice, applied to how we build.", 19)],
], top=Inches(2.2), gap=12)
notes(s, "The portable lessons. One: don't fight the platform — its limits became the design. "
         "Two: the migration lesson is the one I'd tattoo — a schema change isn't done until its "
         "migration ships with it. Three: for a credibility tool, showing the working IS the "
         "product. Four: meta-point — the tool tells users to pick the smallest model that does "
         "the job; same applies to our own engineering. ~55s.")

# 20 · END
s = new_slide()
pill(s, "thank you", top=Inches(2.5))
text(s, MARGIN, Inches(3.2), CONTENT_W, Inches(1.2), [[S("Questions?", sz=52, c=TEXT, b=True)]])
text(s, MARGIN, Inches(4.7), CONTENT_W, Inches(0.5),
     [[S("App · ", 18, c=MUTED), M("ai-carbon-footprint.telus.gizmos.run", sz=16, c=GREEN)]])
text(s, MARGIN, Inches(5.2), CONTENT_W, Inches(0.5),
     [[S("Code · ", 18, c=MUTED), M("github.com/bvillav-a11y/ai-carbon-calculator", sz=16, c=GREEN)]])
notes(s, "Close. Invite questions; offer to walk the live app or the Worker source. Likely Qs: "
         "where the per-token numbers come from (Epoch AI), why embodied is included, how gizmos "
         "D1 differs from raw D1 (the D1Proxy API). ~10s + Q&A.")

# ── save ──
out = os.path.join(os.path.dirname(__file__), "ai-carbon-calculator.pptx")
prs.save(out)
print(f"Wrote {out} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
